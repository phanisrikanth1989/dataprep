#!/usr/bin/env python3
"""Generate DataPrep Architecture Diagram PowerPoint Presentation"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 118, 210)  # Blue

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(title):
    """Add content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 118, 210)

    # Separator line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(200, 200, 200)
    line.line.width = Pt(2)

    return slide

def add_box(slide, left, top, width, height, text, bg_color, text_color=RGBColor(255, 255, 255)):
    """Add a colored box with text"""
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(100, 100, 100)
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)

    p = text_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    return shape

def add_arrow(slide, x1, y1, x2, y2, color=RGBColor(0, 0, 0)):
    """Add arrow between points"""
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(3)

# Slide 1: Title
add_title_slide("DataPrep ETL Engine", "Talend to Python Migration Architecture")

# Slide 2: High-Level Overview
slide = add_content_slide("System Overview")
add_box(slide, 0.5, 1.5, 4, 1.5, "CONVERTER\n(Talend XML to JSON)", RGBColor(66, 165, 245))
add_box(slide, 5.5, 1.5, 4, 1.5, "ENGINE\n(Execute JSON Config)", RGBColor(129, 199, 132))
add_arrow(slide, 4.8, 2.25, 5.2, 2.25, RGBColor(0, 0, 0))

add_box(slide, 0.5, 3.8, 9, 1.2, "Feature Parity Goal: Identical results to Talend for same input", RGBColor(255, 152, 0))
add_box(slide, 0.5, 5.3, 9, 1.5, "Tech Stack: Python 3.12+ | Java 11+ Bridge (Py4J + Arrow) | pandas | pytest", RGBColor(156, 39, 176))

# Slide 3: Converter Pipeline
slide = add_content_slide("Converter Pipeline (12 Steps)")
steps = ["1. Parse XML", "2. Convert\nComponents", "3. Parse Flows", "4. Detect Subjobs"]
for i, step in enumerate(steps):
    add_box(slide, 0.5 + i*2.1, 1.8, 1.9, 1.0, step, RGBColor(66, 165, 245))
    if i < 3:
        add_arrow(slide, 2.3 + i*2.1, 2.3, 2.5 + i*2.1, 2.3)

steps2 = ["5. Map\nTriggers", "6. Convert\nExpressions", "7. Validate\nReferences", "8. Validate\ntMap Rules"]
for i, step in enumerate(steps2):
    add_box(slide, 0.5 + i*2.1, 3.2, 1.9, 1.0, step, RGBColor(100, 181, 246))
    if i < 3:
        add_arrow(slide, 2.3 + i*2.1, 3.7, 2.5 + i*2.1, 3.7)

steps3 = ["9. Validate\nExpressions", "10. Check\nQuality", "11. Final\nValidation", "12. Assemble\nJSON"]
for i, step in enumerate(steps3):
    add_box(slide, 0.5 + i*2.1, 4.6, 1.9, 1.0, step, RGBColor(144, 202, 249))
    if i < 3:
        add_arrow(slide, 2.3 + i*2.1, 5.1, 2.5 + i*2.1, 5.1)

# Slide 4: Data Flow
slide = add_content_slide("Data Flow Through System")
add_box(slide, 1, 1.5, 2.5, 0.8, "Talend\n.item XML", RGBColor(158, 158, 158))
add_arrow(slide, 3.5, 1.9, 3.9, 1.9)
add_box(slide, 4, 1.5, 2.5, 0.8, "XML Parser", RGBColor(66, 165, 245))
add_arrow(slide, 6.5, 1.9, 6.9, 1.9)
add_box(slide, 7, 1.5, 2.5, 0.8, "Converters\n(Registry)", RGBColor(66, 165, 245))

add_box(slide, 1, 2.8, 2.5, 0.8, "Expr. Convert", RGBColor(100, 181, 246))
add_arrow(slide, 3.5, 3.2, 3.9, 3.2)
add_box(slide, 4, 2.8, 2.5, 0.8, "Trigger Map", RGBColor(100, 181, 246))
add_arrow(slide, 6.5, 3.2, 6.9, 3.2)
add_box(slide, 7, 2.8, 2.5, 0.8, "Validate", RGBColor(100, 181, 246))

add_box(slide, 3, 4.2, 4, 1.0, "JSON Job Config", RGBColor(129, 199, 132))
add_arrow(slide, 5, 5.2, 5, 5.6)
add_box(slide, 2.5, 5.7, 5, 0.8, "ETLEngine Execution", RGBColor(76, 175, 80))

# Slide 5: Two-Layer Architecture
slide = add_content_slide("Two-Layer Architecture")
add_box(slide, 0.3, 1.5, 4.2, 5.5, "CONVERTER LAYER\n(talend_to_v1/)\n\nXML Parser\nComponent Converters (80+)\nRegistry Pattern\nExpression Converter\nTrigger Mapper\nValidator\n\nJSON Config Output", RGBColor(66, 165, 245), RGBColor(0, 0, 0))
add_arrow(slide, 4.5, 4.0, 5.5, 4.0, RGBColor(255, 87, 34))
add_box(slide, 5.5, 1.5, 4.2, 5.5, "ENGINE LAYER\n(v1/engine/)\n\nETLEngine Orchestrator\nComponent Registry\nEngine Components (50+)\nBaseComponent ABC\nData Flow Execution\n\nDataFrame Output", RGBColor(129, 199, 132), RGBColor(0, 0, 0))

# Slide 6: Component Registry Pattern
slide = add_content_slide("Component Registry Pattern")
add_box(slide, 0.5, 1.6, 4.3, 1.8, "CONVERTER\n\n@REGISTRY.register()\ndecorator\n\nLookup:\nREGISTRY.get()\n\nABC:\nComponentConverter", RGBColor(66, 165, 245), RGBColor(0, 0, 0))
add_box(slide, 5.2, 1.6, 4.3, 1.8, "ENGINE\n\nManual dict\nregistration\n\nLookup:\nCOMPONENT_REGISTRY.get()\n\nABC:\nBaseComponent", RGBColor(129, 199, 132), RGBColor(0, 0, 0))
add_box(slide, 1.5, 4.0, 7, 0.8, "Both support camelCase (FileInputDelimited) AND Talend names (tFileInputDelimited)", RGBColor(255, 193, 7))
add_box(slide, 0.5, 5.1, 9, 1.8, "Why Registry Pattern?\n✓ Automatic component discovery  ✓ Flexible lookup by name\n✓ Easy to add new components  ✓ Consistent discipline", RGBColor(156, 39, 176))

# Slide 7: BaseComponent Pattern
slide = add_content_slide("BaseComponent Execution Flow")
add_box(slide, 0.5, 1.5, 1.8, 0.7, "execute()", RGBColor(66, 165, 245), RGBColor(0, 0, 0))
add_arrow(slide, 1.4, 2.2, 1.4, 2.4)
add_box(slide, 0.3, 2.5, 2.2, 1.2, "1. Resolve\n{{java}}\nexpressions", RGBColor(100, 181, 246), RGBColor(0, 0, 0))
add_arrow(slide, 1.4, 3.7, 1.4, 3.9)
add_box(slide, 0.3, 4.0, 2.2, 1.2, "2. Resolve\n${context.var}\nvariables", RGBColor(100, 181, 246), RGBColor(0, 0, 0))
add_arrow(slide, 1.4, 5.2, 1.4, 5.4)
add_box(slide, 0.3, 5.5, 2.2, 1.0, "3. Auto-select\nExecution Mode", RGBColor(100, 181, 246), RGBColor(0, 0, 0))

add_box(slide, 3.3, 2.5, 1.8, 1.2, "4. Call\n_process()\n(abstract)", RGBColor(129, 199, 132), RGBColor(0, 0, 0))
add_arrow(slide, 4.2, 3.7, 4.2, 3.9)
add_box(slide, 3.3, 4.0, 1.8, 1.2, "5. Return\n{main, reject,\nstats}", RGBColor(129, 199, 132), RGBColor(0, 0, 0))
add_arrow(slide, 4.2, 5.2, 4.2, 5.4)
add_box(slide, 3.3, 5.5, 1.8, 1.0, "6. Update\nGlobalMap\nStats", RGBColor(129, 199, 132), RGBColor(0, 0, 0))

add_box(slide, 5.5, 1.5, 4, 5.8, "Three Execution Modes\n\n• BATCH: Entire\nDataFrame at once\n\n• STREAMING: Row-by-row\n(memory efficient)\n\n• HYBRID: Batch +\niterator support\n\nKey Methods:\n_update_global_map()\n_update_stats()\nvalidate_schema()\n_process() [abstract]", RGBColor(255, 152, 0), RGBColor(0, 0, 0))

# Slide 8: Java Bridge Architecture
slide = add_content_slide("Java Bridge Architecture (Py4J + Arrow)")
add_box(slide, 0.3, 1.5, 2.8, 1.5, "PYTHON\nEngine\n\nJavaBridgeManager\nBridge.py", RGBColor(66, 165, 245), RGBColor(0, 0, 0))
add_arrow(slide, 3.1, 2.25, 3.5, 2.25, RGBColor(255, 87, 34))
add_box(slide, 3.6, 1.5, 2.8, 1.5, "JAVA (JVM)\nBridge Server\n\nJavaBridge.java\nRowWrapper.java", RGBColor(129, 199, 132), RGBColor(0, 0, 0))
add_box(slide, 0.3, 3.3, 9.4, 0.8, "Transport: Py4J (RPC) + Apache Arrow (binary row data)", RGBColor(156, 39, 176))
add_box(slide, 0.3, 4.3, 4.5, 2.5, "Python to Java\n\n1. Sync context +\n   globalMap\n\n2. Batch {{java}}\n   expressions\n\n3. Send via Arrow\n\n4. Execute Groovy", RGBColor(100, 181, 246), RGBColor(0, 0, 0))
add_box(slide, 5.2, 4.3, 4.5, 2.5, "Java to Python\n\n1. Execute tMap\n   transformations\n\n2. Serialize via\n   Apache Arrow\n\n3. Return results\n\n4. Sync back state", RGBColor(144, 202, 249), RGBColor(0, 0, 0))

# Slide 9: Global State Management
slide = add_content_slide("Global State & Context Management")
add_box(slide, 0.5, 1.5, 4, 2.3, "GlobalMap\n\nTalend-compatible\nkey-value store\n\n• NB_LINE\n• NB_LINE_OK\n• NB_LINE_REJECT\n• Custom vars", RGBColor(66, 165, 245), RGBColor(0, 0, 0))
add_arrow(slide, 4.7, 2.6, 5.3, 2.6)
add_box(slide, 5.5, 1.5, 4, 2.3, "ContextManager\n\nResolve patterns\n${context.var}\n\n• Load files\n• Type coercion\n• Defaults", RGBColor(129, 199, 132), RGBColor(0, 0, 0))
add_box(slide, 0.5, 4.1, 9, 1.2, "TriggerManager: Orchestrates flow transitions\nOnComponentOk | OnSubjobOk | RunIf", RGBColor(255, 152, 0))
add_box(slide, 0.5, 5.5, 9, 1.5, "Data Flow Integration\nGlobalMap ←→ Java Bridge ←→ Components → Output Router → data_flows", RGBColor(156, 39, 176))

# Slide 10: Directory Structure
slide = add_content_slide("Directory Structure")
add_box(slide, 0.3, 1.5, 4.5, 1.2, "src/converters/\ntalend_to_v1/\n\n80+ Converter\nComponents", RGBColor(66, 165, 245), RGBColor(0, 0, 0))
add_box(slide, 0.3, 2.9, 4.5, 1.2, "src/v1/engine/\n\n50+ Engine\nComponents", RGBColor(129, 199, 132), RGBColor(0, 0, 0))
add_box(slide, 0.3, 4.3, 4.5, 1.2, "src/v1/\njava_bridge/\n\nPython + Java\nserver", RGBColor(100, 181, 246), RGBColor(0, 0, 0))
add_box(slide, 5.1, 1.5, 4.5, 1.2, "tests/converters/\ntalend_to_v1/\n\nTest coverage", RGBColor(255, 152, 0), RGBColor(0, 0, 0))
add_box(slide, 5.1, 2.9, 4.5, 1.2, "tests/v1/engine/\n\n258 test files\n95% gate", RGBColor(255, 152, 0), RGBColor(0, 0, 0))
add_box(slide, 5.1, 4.3, 4.5, 1.2, "tests/fixtures/\n\nTest data +\nconfigs", RGBColor(255, 152, 0), RGBColor(0, 0, 0))
add_box(slide, 0.5, 5.8, 9, 1.0, "Categories: aggregate | context | control | database | file | iterate | transform", RGBColor(156, 39, 176))

# Slide 11: Error Handling
slide = add_content_slide("Error Handling Hierarchy")
add_box(slide, 1.5, 1.5, 7, 1.0, "ETLError (Root)", RGBColor(244, 67, 54))
errors = [("ConfigurationError", 0), ("DataValidationError", 2), ("ComponentExecutionError", 4),
          ("FileOperationError", 6), ("JavaBridgeError", 1), ("ExpressionError", 3), ("SchemaError", 5)]
for err, col in errors:
    x = 0.7 + (col % 3) * 3.1
    y = 2.8 + (col // 3) * 1.2
    add_box(slide, x, y, 2.8, 1.0, err, RGBColor(255, 152, 0), RGBColor(0, 0, 0))

add_box(slide, 0.5, 5.5, 9, 1.7, "Engine Behavior:\n• Caught in _execute_component()  • Component marked failed\n• Error recorded in execution_stats  • Execution continues (unless die_on_error)", RGBColor(76, 175, 80))

# Slide 12: Metrics
slide = add_content_slide("Project Metrics & Coverage Gate")
add_box(slide, 0.5, 1.5, 2.3, 1.5, "Code Stats\n\n211 Source\n\n258 Tests\n\n469 Total", RGBColor(66, 165, 245), RGBColor(0, 0, 0))
add_box(slide, 3.2, 1.5, 2.3, 1.5, "Components\n\n80+ Converter\n\n50+ Engine\n\n130+ Total", RGBColor(129, 199, 132), RGBColor(0, 0, 0))
add_box(slide, 5.9, 1.5, 3.6, 1.5, "Coverage\n\n95% per-module\nline coverage\n\nJava included", RGBColor(255, 152, 0), RGBColor(0, 0, 0))
add_box(slide, 0.5, 3.3, 9, 0.9, "Test Markers: @unit | @integration | @java | @oracle (opt-in) | @slow", RGBColor(100, 181, 246))
add_box(slide, 0.5, 4.4, 9, 2.6, "Dependencies (Modular)\nCore: pandas, numpy\nJava: pyarrow, py4j | Excel: openpyxl | Oracle: oracledb | XML: lxml\nYAML: PyYAML | JSON: jsonpath-ng | API: fastapi | Dev: pytest\n\nPython 3.12+ | JVM 11+ | Maven 3.x", RGBColor(156, 39, 176))

# Slide 13: Conventions
slide = add_content_slide("Coding Conventions")
add_box(slide, 0.5, 1.5, 4.4, 3.5, "Naming\n\nModules:\nsnake_case.py\n\nClasses:\nPascalCase\n\nFunctions:\nsnake_case\n\nConstants:\nUPPER_SNAKE_CASE\n\nExceptions:\n*Error", RGBColor(66, 165, 245), RGBColor(0, 0, 0))
add_box(slide, 5.2, 1.5, 4.4, 3.5, "Standards\n\n4-space indent\n\nDouble quotes\n\nType hints\n\nASCII-only logs\n\nModule loggers\n\nPublic docstrings", RGBColor(129, 199, 132), RGBColor(0, 0, 0))
add_box(slide, 0.5, 5.2, 9, 1.8, "Imports\nRelative in packages: from ..base_component import BaseComponent\nAbsolute in tests: from src.v1.engine import ETLEngine", RGBColor(156, 39, 176))

# Slide 14: Quick Reference
slide = add_content_slide("Quick Reference - Key Files")
add_box(slide, 0.3, 1.5, 4.4, 2.5, "Converter Files\n\nconverter.py:460\n  Main entry point\n\nxml_parser.py\n  XML to TalendJob\n\ncomponents/base.py\n  ABC class\n\ncomponents/registry.py\n  @register decorator", RGBColor(66, 165, 245), RGBColor(0, 0, 0))
add_box(slide, 5.1, 1.5, 4.4, 2.5, "Engine Files\n\nengine.py:860\n  Main entry point\n\nbase_component.py\n  ABC class\n\ncomponent_registry.py\n  Registry dict\n\nexecutor.py\n  Execution loop", RGBColor(129, 199, 132), RGBColor(0, 0, 0))
add_box(slide, 0.3, 4.2, 4.4, 2.5, "Services\n\ncontext_manager.py\n  ${context.var}\n\nglobal_map.py\n  Statistics\n\njava_bridge_manager.py\n  Py4J lifecycle", RGBColor(100, 181, 246), RGBColor(0, 0, 0))
add_box(slide, 5.1, 4.2, 4.4, 2.5, "Support\n\nexceptions.py\n  Error hierarchy\n\nexecution_plan.py\n  DAG planning\n\noutput_router.py\n  Flow management", RGBColor(144, 202, 249), RGBColor(0, 0, 0))

# Save presentation
import os
output_path = os.path.join(os.getcwd(), "DataPrep_Architecture_Diagrams.pptx")
prs.save(output_path)
print("[OK] PowerPoint created: " + output_path)
print("[OK] 14 professional slides with clear architecture diagrams")
print("[OK] Total: 469 files analyzed | 130+ components documented")
